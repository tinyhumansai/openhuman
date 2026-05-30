"""Generate a PowerPoint (.pptx) file from a structured slide spec on stdin.

Contract (kept in sync with `types.rs::GeneratePresentationInput` on
the Rust side):

    stdin  : UTF-8 JSON object
              { "title": "...",
                "author": "..." | null,
                "theme": "default" | "minimal" | "dark" | null,
                "slides": [
                  { "title": "...",
                    "body": "..." | null,
                    "bullets": [...] | null,
                    "speaker_notes": "..." | null }
                ] }
    argv   : --output <ABSOLUTE_PATH_TO_PPTX>
    stdout : JSON `{"ok": true, "slide_count": N}` on success
    stderr : human-readable error text on failure
    exit   : 0 on success, 2 on input/runtime error, 3 on python-pptx
             error

This script must remain a single self-contained file. No relative
imports, no `eval`, no network. It is shipped into the binary via
`include_str!` and materialised to a tempfile per invocation by the
caller in `script.rs`.
"""

from __future__ import annotations

import argparse
import json
import os
import sys


def _exit_error(message: str, exit_code: int = 2) -> None:
    sys.stderr.write(message.rstrip() + "\n")
    sys.exit(exit_code)


def _coerce_str(value, field: str, allow_empty: bool = True) -> str:
    if value is None:
        return ""
    if not isinstance(value, str):
        _exit_error(f"field {field!r} must be a string, got {type(value).__name__}")
    if not allow_empty and not value.strip():
        _exit_error(f"field {field!r} must be a non-empty string")
    return value


def _coerce_list_of_str(value, field: str) -> list[str]:
    if value is None:
        return []
    if not isinstance(value, list):
        _exit_error(f"field {field!r} must be a list of strings")
    out = []
    for i, item in enumerate(value):
        if not isinstance(item, str):
            _exit_error(f"field {field!r}[{i}] must be a string")
        out.append(item)
    return out


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a .pptx from a JSON slide spec")
    parser.add_argument(
        "--output",
        required=True,
        help="absolute path where the .pptx should be written",
    )
    args = parser.parse_args()

    # Defence-in-depth: the Rust caller always supplies an absolute
    # artifact path, but a future regression there must not silently
    # redirect output to a cwd-relative location.
    if not os.path.isabs(args.output):
        _exit_error("argument '--output' must be an absolute path")
        return

    try:
        raw = sys.stdin.read()
    except OSError as err:
        _exit_error(f"failed to read stdin: {err}")
        return

    if not raw.strip():
        _exit_error("stdin payload was empty; expected a JSON slide spec")
        return

    try:
        spec = json.loads(raw)
    except json.JSONDecodeError as err:
        _exit_error(f"stdin payload is not valid JSON: {err}")
        return

    if not isinstance(spec, dict):
        _exit_error("stdin payload must be a JSON object")
        return

    title = _coerce_str(spec.get("title"), "title", allow_empty=False)
    author = _coerce_str(spec.get("author"), "author")
    # theme is accepted but currently informational only; the
    # python-pptx default template is used regardless. Future
    # iteration can swap layouts based on `theme`.
    _ = _coerce_str(spec.get("theme"), "theme")
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        _exit_error("field 'slides' must be a non-empty list")
        return

    try:
        from pptx import Presentation
    except ImportError as err:
        _exit_error(
            "python-pptx is not installed in this interpreter; expected via runtime_python venv. "
            f"underlying error: {err}",
            exit_code=2,
        )
        return

    try:
        prs = Presentation()

        # Title slide.
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if author and len(slide.placeholders) > 1:
            slide.placeholders[1].text = author

        # Content slides.
        content_layout = prs.slide_layouts[1]
        for i, raw_slide in enumerate(slides):
            if not isinstance(raw_slide, dict):
                _exit_error(f"slides[{i}] must be a JSON object")
                return
            slide_title = _coerce_str(raw_slide.get("title"), f"slides[{i}].title")
            body = _coerce_str(raw_slide.get("body"), f"slides[{i}].body")
            bullets = _coerce_list_of_str(
                raw_slide.get("bullets"), f"slides[{i}].bullets"
            )
            notes = _coerce_str(raw_slide.get("speaker_notes"), f"slides[{i}].speaker_notes")

            slide = prs.slides.add_slide(content_layout)
            if slide_title:
                slide.shapes.title.text = slide_title

            # python-pptx's default "Title and Content" layout exposes
            # the body placeholder at index 1. Defensive lookup keeps
            # the script robust against custom templates where the
            # ordering shifts.
            body_placeholder = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body_placeholder = ph
                    break

            if body_placeholder is not None and (body or bullets):
                text_frame = body_placeholder.text_frame
                first_line_placed = False
                if body.strip():
                    text_frame.text = body
                    first_line_placed = True
                for bullet in bullets:
                    if first_line_placed:
                        para = text_frame.add_paragraph()
                    else:
                        para = text_frame.paragraphs[0]
                        first_line_placed = True
                    para.text = bullet
                    para.level = 0

            if notes.strip():
                slide.notes_slide.notes_text_frame.text = notes

        prs.save(args.output)
    except Exception as err:  # noqa: BLE001 — python-pptx raises many concrete types
        _exit_error(f"python-pptx generation failed: {err}", exit_code=3)
        return

    sys.stdout.write(json.dumps({"ok": True, "slide_count": len(slides)}))
    sys.stdout.write("\n")
    sys.exit(0)


if __name__ == "__main__":
    main()
